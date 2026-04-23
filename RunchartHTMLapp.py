# runcharts_streamlit_premium.py
# PPTX + HTML Dashboard (SAFE UPGRADE - LOGIC PRESERVED)

import streamlit as st
import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE

import plotly.graph_objects as go
import re
import difflib
import io
import os
from datetime import datetime

# ---------------------------
# CONFIG (UNCHANGED)
# ---------------------------
NUM_POINTS = 18
CENTER_LINE_METHOD = "median"
ASTRO_THRESHOLD = 10

NON_DATA_COLS_LOWER = {"department", "indicator", "target", "benchmark/ category", "benchmark", "frequency"}

TITLE_BOX = (0.6, 0.4, 9.0, 0.9)
CHART_BOX = (0.6, 1.4, 9.0, 4.3)
NOTES_BOX = (0.6, 6.0, 9.0, 1.3)

LOG_DIR = "logs"
UPLOAD_DIR = os.path.join(LOG_DIR, "uploads")
LOG_FILE = os.path.join(LOG_DIR, "activity_log.csv")

os.makedirs(UPLOAD_DIR, exist_ok=True)

# ---------------------------
# LOGIN
# ---------------------------
PREMIUM_LOGIN = "Pakistan@1947"
PREMIUM_PASSWORD = "Pakistan@1947"

# ---------------------------
# CLEAN FUNCTIONS (UNCHANGED)
# ---------------------------
def clean_text_for_match(val):
    if pd.isna(val):
        return ""
    s = str(val)
    s = re.sub(r"[\u200B\u200C\u200D\uFEFF]", "", s)
    s = re.sub(r"\s+", " ", s)
    return s.strip().lower()

def pretty_label(col):
    try:
        parsed = pd.to_datetime(col, errors="coerce")
        if not pd.isna(parsed):
            return parsed.strftime("%b-%y")
    except:
        pass
    return str(col)

def get_center_line(values, method="median"):
    arr = np.array(values, dtype=float)
    arr = arr[~np.isnan(arr)]
    if arr.size == 0:
        return np.nan
    return float(np.nanmedian(arr)) if method == "median" else float(np.nanmean(arr))

# ---------------------------
# DETECTION LOGIC (UNCHANGED)
# ---------------------------
def detect_shift(series, center, min_run=6):
    signs = []
    for v in series:
        if pd.isna(v):
            signs.append(0)
        elif v > center:
            signs.append(1)
        elif v < center:
            signs.append(-1)
        else:
            signs.append(0)

    shifts = []
    i, n = 0, len(signs)

    while i < n:
        if signs[i] == 0:
            i += 1
            continue

        curr = signs[i]
        j = i + 1
        count = 1

        while j < n and signs[j] == curr:
            count += 1
            j += 1

        if count >= min_run:
            shifts.append((i, j - 1, curr))

        i = j

    return shifts

def detect_trend(series, min_run=5):
    trends = []
    vals = [v for v in series if pd.notna(v)]

    for i in range(len(vals) - min_run):
        window = vals[i:i+min_run]
        if all(window[j] < window[j+1] for j in range(len(window)-1)):
            trends.append((i, i+min_run-1, 1))
        if all(window[j] > window[j+1] for j in range(len(window)-1)):
            trends.append((i, i+min_run-1, -1))

    return trends

def detect_astronomical(series, median_val, threshold=10):
    ast = []
    for i in range(1, len(series)):
        if pd.notna(series[i]) and pd.notna(series[i-1]):
            if abs(series[i] - series[i-1]) >= threshold:
                ast.append(i)
    return ast

# ---------------------------
# PPTX GENERATION (UNCHANGED CORE)
# ---------------------------
def create_presentation_for_department(slide_infos, dept_display_name):
    prs = Presentation()

    for info in slide_infos:
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        title = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(9), Inches(0.6))
        title.text_frame.text = info["indicator_name"]

        labels = info["labels"]
        series = info["series"]

        chart_data = CategoryChartData()
        chart_data.categories = labels
        chart_data.add_series("Value", series)

        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE_MARKERS,
            Inches(0.6), Inches(1.2), Inches(9), Inches(4),
            chart_data
        ).chart

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output

# ---------------------------
# LOGIC ENGINE
# ---------------------------
def build_analysis(df, dept_col, ind_col, data_cols, dept_clean):
    slide_infos = []

    df_dept = df[df["_dept_clean"] == dept_clean]

    for _, row in df_dept.iterrows():

        series = []
        labels = []

        for c in data_cols:
            labels.append(pretty_label(c))
            val = pd.to_numeric(row.get(c, np.nan), errors="coerce")
            series.append(val)

        center = get_center_line(series)

        shifts = detect_shift(series, center)
        trends = detect_trend(series)
        astro = detect_astronomical(series, center)

        imp = []

        for s in shifts:
            imp.append(f"SHIFT {s}")

        for t in trends:
            imp.append(f"TREND {t}")

        for a in astro:
            imp.append(f"ASTRO {a}")

        slide_infos.append({
            "indicator_name": str(row[ind_col]),
            "labels": labels,
            "series": series,
            "analysis": imp
        })

    return slide_infos

# ---------------------------
# PLOTLY CHART (NEW)
# ---------------------------
def plot_chart(labels, series, median):
    fig = go.Figure()

    fig.add_trace(go.Scatter(
        x=labels,
        y=series,
        mode="lines+markers",
        name="Value"
    ))

    fig.add_hline(y=median, line_dash="dash", line_color="gray")

    fig.update_layout(height=450, template="simple_white")
    return fig

# ---------------------------
# STREAMLIT UI STATE
# ---------------------------
st.set_page_config(layout="wide")
st.title("📊 RunCharts Premium Dashboard")

if "page" not in st.session_state:
    st.session_state.page = "upload"

if "df" not in st.session_state:
    st.session_state.df = None

# ---------------------------
# PAGE 1 - UPLOAD
# ---------------------------
if st.session_state.page == "upload":

    file = st.file_uploader("Upload Excel")

    if file:
        df = pd.read_excel(file)
        st.session_state.df = df

        st.success("File Loaded")

        if st.button("Go to Dashboard"):
            st.session_state.page = "dashboard"
            st.rerun()

# ---------------------------
# PAGE 2 - DASHBOARD (DEPARTMENTS)
# ---------------------------
elif st.session_state.page == "dashboard":

    df = st.session_state.df

    dept_col = st.selectbox("Department Column", df.columns)
    ind_col = st.selectbox("Indicator Column", df.columns)

    df["_dept_clean"] = df[dept_col].astype(str).apply(clean_text_for_match)

    departments = df["_dept_clean"].unique()

    st.subheader("Departments")

    for d in departments:
        if st.button(f"📁 {d}"):
            st.session_state.selected_dept = d
            st.session_state.page = "indicators"
            st.rerun()

    if st.button("⬅ Back"):
        st.session_state.page = "upload"
        st.rerun()

# ---------------------------
# PAGE 3 - INDICATORS
# ---------------------------
elif st.session_state.page == "indicators":

    df = st.session_state.df
    dept = st.session_state.selected_dept

    df = df[df["_dept_clean"] == dept]

    st.subheader(f"Indicators - {dept}")

    for i, row in df.iterrows():
        if st.button(str(row.iloc[1])):

            st.session_state.selected_indicator = i
            st.session_state.page = "chart"
            st.rerun()

    if st.button("🏠 Main Menu"):
        st.session_state.page = "dashboard"
        st.rerun()

# ---------------------------
# PAGE 4 - CHART
# ---------------------------
elif st.session_state.page == "chart":

    df = st.session_state.df
    row = df.iloc[st.session_state.selected_indicator]

    labels = list(df.columns[2:])
    series = [pd.to_numeric(row[c], errors="coerce") for c in labels]

    median = np.nanmedian(series)

    st.subheader(row.iloc[1])

    st.plotly_chart(plot_chart(labels, series, median), use_container_width=True)

    st.write("Analysis:", series)

    if st.button("⬅ Back"):
        st.session_state.page = "indicators"
        st.rerun()

    if st.button("🏠 Main Menu"):
        st.session_state.page = "dashboard"
        st.rerun()
